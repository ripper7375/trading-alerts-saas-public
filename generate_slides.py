from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Trading Alerts SaaS"
    subtitle.text = "Types Definition & Tier System Overview\n(Part 03 & Part 04 Architecture)"

    # Slide 2: System Overview
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "System Overview"
    tf = body_shape.text_frame
    tf.text = "Comprehensive Type Safety: No 'any' types, strict null checking, comprehensive enums."
    p = tf.add_paragraph()
    p.text = "Unified Tier Architecture: Streamlined from 3 to 2 tiers (FREE & PRO)."
    p = tf.add_paragraph()
    p.text = "Advanced Indicator System: 57-column database schema with dynamic filtering."
    p = tf.add_paragraph()
    p.text = "Seamless Integrations: Multi-provider payments (Stripe, dLocal), RiseWorks payouts, TOTP 2FA."

    # Slide 3: Tier System - FREE vs PRO
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Tier System Features: FREE vs PRO"
    tf = body_shape.text_frame
    
    tf.text = "FREE Tier ($0/month):"
    p = tf.add_paragraph()
    p.text = "5 Symbols & 3 Timeframes (15 chart combinations)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Limits: 5 alerts, 1 watchlist (5 items)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Access: 24 DB columns (8 system + 16 basic indicators)."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "PRO Tier ($29/month, 7-Day Trial):"
    p = tf.add_paragraph()
    p.text = "15 Symbols & 9 Timeframes (135 chart combinations)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Limits: 20 alerts, 5 watchlists (50 items)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Access: Full 57 DB columns (all 8 advanced indicators)."
    p.level = 1

    # Slide 4: Indicator & Data Architecture
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Indicator & Data Architecture"
    tf = body_shape.text_frame
    
    tf.text = "57-Column Database Schema:"
    p = tf.add_paragraph()
    p.text = "System Columns (8): OHLCV & metadata."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "FREE Indicators (16): Fractal Horizontal, Fractal Diagonal."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "PRO Indicators (33): Moving Averages, Body Momentum, Heiken Ashi, Keltner Channels, Support/Resistance, ZigZag."
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "Tier-Aware API Responses:"
    p = tf.add_paragraph()
    p.text = "Data dynamically filtered based on user tier via filterDataByTier()."
    p.level = 1

    # Slide 5: Tier Validation & Access Control
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Tier Validation & Access Control"
    tf = body_shape.text_frame
    
    tf.text = "Strict Server-Side Validation Functions:"
    p = tf.add_paragraph()
    p.text = "Symbol/Timeframe validation (e.g., canAccessSymbol, validateChartAccess)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Usage limits enforcement (e.g., canCreateAlert, canCreateWatchlist)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Indicator/Column gating (e.g., canAccessColumn, canAccessIndicator)."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Performance Metrics:"
    p = tf.add_paragraph()
    p.text = "API Requests: 60/hr (FREE), 300/hr (PRO)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Chart Updates: Real-time OHLCV via WebSocket for BOTH tiers."
    p.level = 1

    # Slide 6: Core Domain Types Definition
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Core Domain Types & Security"
    tf = body_shape.text_frame
    
    tf.text = "Users & Authentication:"
    p = tf.add_paragraph()
    p.text = "Extended NextAuth Sessions (tier, role, affiliate status)."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "TOTP 2FA fields and trial management tracking."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Payments & Operations:"
    p = tf.add_paragraph()
    p.text = "dLocal (8 emerging market countries) & Stripe standard subscriptions."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "RiseWorks payouts, KYC, Batch payments tracking."
    p.level = 1

    # Slide 7: Development Ecosystem
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    title_shape.text = "Development & Tooling Ecosystem"
    tf = body_shape.text_frame
    
    tf.text = "Prisma Stubs (types/prisma-stubs.d.ts):"
    p = tf.add_paragraph()
    p.text = "Supports offline dev and environments without Prisma client."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Full type coverage for models and Enums."
    p.level = 1

    p = tf.add_paragraph()
    p.text = "Documentation:"
    p = tf.add_paragraph()
    p.text = "100% typed definitions exported via types/index.ts."
    p.level = 1
    p = tf.add_paragraph()
    p.text = "Fully mapped in OpenAPI specifications."
    p.level = 1

    prs.save('Trading_Alerts_SaaS_Types_And_Tiers.pptx')
    print("Successfully generated Trading_Alerts_SaaS_Types_And_Tiers.pptx")

if __name__ == '__main__':
    create_presentation()
